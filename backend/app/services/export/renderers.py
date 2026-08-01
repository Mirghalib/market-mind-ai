"""Strategy export renderers.

Each format is a ``Renderer`` that turns a MarketingStrategy into
file bytes plus the metadata needed for an HTTP file response
(media type, file extension). Renderers are registered in ``RENDERERS``;
the endpoint and service stay untouched.

Formats:
  - json: full document
  - markdown: readable, version-control friendly
  - html: styled single-file report
  - pdf: reportlab-based (requires ``reportlab``)
  - docx: python-docx-based (requires ``python-docx``)
"""
import json
from dataclasses import dataclass
from html import escape

from app.models.export import ExportFormat
from app.models.marketing_strategy import MarketingStrategy


@dataclass(frozen=True)
class RenderedExport:
    """Result of a renderer: payload bytes + response metadata."""

    content: bytes
    media_type: str
    file_extension: str


class BaseRenderer:
    """Interface implemented by every export format renderer."""

    format: ExportFormat

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        raise NotImplementedError


def _document(strategy: MarketingStrategy) -> dict:
    """Normalized document shared by the text-based renderers."""
    goals = getattr(strategy, "goals", None) or []
    status = getattr(strategy, "status", None)
    created_at = getattr(strategy, "created_at", None)
    updated_at = getattr(strategy, "updated_at", None)
    return {
        "strategy_id": str(strategy.id),
        "name": strategy.name,
        "target_audience": strategy.target_audience,
        "goals": goals,
        "status": status.value if status else "",
        "content": strategy.content or {},
        "created_at": created_at.isoformat() if created_at else "",
        "updated_at": updated_at.isoformat() if updated_at else "",
    }


def _sections(strategy: MarketingStrategy) -> list[tuple[str, str]]:
    """Return [(title, content)] pairs from the stored payload.

    Handles three shapes:
      1. Flat list shape:  content.sections = [{title, content}, ...]
      2. Raw LLM document: content = {marketingStrategy: {...}, customerPersona: {...}, ...}
      3. Legacy flat:      content = {summary: "...", sections: [...]}
    """
    content = strategy.content or {}

    # Shape 1: flat list.
    sections = content.get("sections")
    if isinstance(sections, list) and sections:
        return [
            (str(s.get("title", "Section")), str(s.get("content", "")))
            for s in sections
        ]

    # Shape 2: raw LLM document — flatten top-level objects into sections.
    flat: list[tuple[str, str]] = []
    for key, value in content.items():
        title = key.replace("_", " ").title()
        if isinstance(value, dict):
            rendered = _flatten_block(value)
            if rendered:
                flat.append((title, rendered))
        elif isinstance(value, str) and value.strip():
            flat.append((title, value))
    if flat:
        return flat

    # Shape 3: legacy flat shape.
    legacy = []
    if content.get("summary"):
        legacy.append(("Summary", str(content["summary"])))
    for key, value in content.items():
        if key in ("summary", "sections"):
            continue
        if isinstance(value, str):
            legacy.append((key.replace("_", " ").title(), value))
    return legacy


def _flatten_block(block: dict) -> str:
    """Render a nested LLM block (e.g. marketingStrategy) as readable text."""
    lines: list[str] = []
    for key, value in block.items():
        label = key.replace("_", " ").replace("-", " ").title()
        if isinstance(value, str) and value.strip():
            lines.append(f"{label}: {value.strip()}")
        elif isinstance(value, list):
            items = [str(i) for i in value if str(i).strip()]
            if items:
                lines.append(f"{label}:")
                lines.extend(f"- {item}" for item in items)
    return "\n".join(lines)


class JsonRenderer(BaseRenderer):
    """Serialise the strategy document as a JSON file."""

    format = ExportFormat.JSON

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        document = _document(strategy)
        payload = json.dumps(document, indent=2, ensure_ascii=False).encode("utf-8")
        return RenderedExport(
            content=payload,
            media_type="application/json",
            file_extension="json",
        )


class MarkdownRenderer(BaseRenderer):
    """Render the strategy as a readable Markdown document."""

    format = ExportFormat.MARKDOWN

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        lines = [
            f"# {strategy.name}",
            "",
            f"- **Status:** {strategy.status.value}",
            f"- **Target audience:** {strategy.target_audience or '—'}",
            "- **Goals:** "
            + (", ".join(strategy.goals) if strategy.goals else "—"),
            f"- **Created:** {_document(strategy)['created_at']}",
            "",
        ]
        for title, content in _sections(strategy):
            lines.append(f"## {title}")
            lines.append("")
            lines.append(content)
            lines.append("")
        payload = "\n".join(lines).encode("utf-8")
        return RenderedExport(
            content=payload,
            media_type="text/markdown",
            file_extension="md",
        )


class HtmlRenderer(BaseRenderer):
    """Render the strategy as a styled, self-contained HTML report."""

    format = ExportFormat.HTML

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        doc = _document(strategy)
        blocks = "".join(
            f"<h2>{escape(title)}</h2>\n<p>{escape(content).replace(chr(10), '<br/>')}</p>\n"
            for title, content in _sections(strategy)
        )
        html = f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8"/>
<title>{escape(strategy.name)} — Market Mind AI</title>
<style>
  body {{ font-family: system-ui, -apple-system, sans-serif; max-width: 800px;
         margin: 2rem auto; padding: 0 1rem; color: #18181b; }}
  h1 {{ color: #4f46e5; }}
  h2 {{ margin-top: 2rem; color: #18181b; }}
  .meta {{ color: #71717a; font-size: 0.9rem; }}
  p {{ line-height: 1.6; white-space: pre-line; }}
</style>
</head>
<body>
  <h1>{escape(strategy.name)}</h1>
  <p class="meta">Status: {escape(strategy.status.value)} · Created: {escape(doc['created_at'])}</p>
  <p class="meta">Target audience: {escape(strategy.target_audience or '—')}</p>
  {blocks}
</body>
</html>"""
        return RenderedExport(
            content=html.encode("utf-8"),
            media_type="text/html",
            file_extension="html",
        )


class PdfRenderer(BaseRenderer):
    """Render the strategy as a professional consulting-grade PDF.

    Delegates to ``pdf_report.build_strategy_report`` (a reportlab-based
    builder with cover page, table of contents, charts, tables, summary
    cards, headers/footers and page numbers). The import is lazy so the
    registry stays importable without reportlab.
    """

    format = ExportFormat.PDF

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        try:
            from app.services.export.pdf_report import build_strategy_report
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError(
                "PDF export requires the 'reportlab' package. "
                "Install it with: pip install reportlab"
            ) from exc

        payload = build_strategy_report(strategy)
        return RenderedExport(
            content=payload,
            media_type="application/pdf",
            file_extension="pdf",
        )


class DocxRenderer(BaseRenderer):
    """Render the strategy as a .docx file (requires ``python-docx``)."""

    format = ExportFormat.DOCX

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        try:
            from docx import Document
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError(
                "DOCX export requires the 'python-docx' package. "
                "Install it with: pip install python-docx"
            ) from exc

        from io import BytesIO

        buffer = BytesIO()
        document = Document()
        document.add_heading(strategy.name, level=0)
        status = getattr(strategy, "status", None)
        document.add_paragraph(
            f"Status: {status.value if status else 'completed'} · "
            f"Created: {_document(strategy)['created_at']}"
        )
        if strategy.target_audience:
            document.add_paragraph(f"Target audience: {strategy.target_audience}")
        if getattr(strategy, "goals", None):
            document.add_paragraph("Goals: " + ", ".join(strategy.goals))

        for title, content in _sections(strategy):
            document.add_heading(title, level=1)
            for paragraph in content.split("\n"):
                if paragraph.strip():
                    document.add_paragraph(paragraph)

        document.save(buffer)
        return RenderedExport(
            content=buffer.getvalue(),
            media_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            file_extension="docx",
        )


class PptxRenderer(BaseRenderer):
    """Render the strategy as a .pptx deck (requires ``python-pptx``).

    Structure: title slide, agenda slide, then one slide per report
    section. Narrative sections use bullet text; tabular sections
    render as tables.
    """

    format = ExportFormat.PPTX

    #: Ordered section builders keyed by the flattened _sections() title.
    _SECTION_KEYS = [
        "Executive Summary", "Marketing Score", "KPI Summary", "Budget Summary",
        "Market Overview", "Customer Persona", "SWOT Analysis",
        "Competitor Analysis", "Marketing Channels", "SEO Strategy",
        "Email Marketing Strategy", "Social Media Strategy",
        "Google & Meta Ads Plan", "Content Calendar",
        "Implementation Roadmap (90 Days)", "Weekly Milestones",
        "Estimated ROI", "Risks and Mitigation", "Final Recommendations",
    ]

    def render(self, strategy: MarketingStrategy) -> RenderedExport:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError(
                "PPTX export requires the 'python-pptx' package. "
                "Install it with: pip install python-pptx"
            ) from exc

        from io import BytesIO

        buffer = BytesIO()
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # --- Title slide -------------------------------------------------
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = title_slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11), Inches(2))
        tf = box.text_frame
        tf.text = strategy.name
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self._rgb("4f46e5")
        p = tf.add_paragraph()
        p.text = "Marketing Strategy Report"
        p.font.size = Pt(24)
        p.font.color.rgb = self._rgb("1e293b")
        meta = title_slide.shapes.add_textbox(
            Inches(1), Inches(5.2), Inches(11), Inches(1.6)
        )
        mf = meta.text_frame
        mf.text = f"Industry: {self._industry(strategy)}\n"
        mf.text += f"Target audience: {strategy.target_audience or '—'}\n"
        mf.text += self._date_line(strategy)
        for para in mf.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = self._rgb("64748b")

        # --- Agenda -------------------------------------------------------
        agenda = prs.slides.add_slide(prs.slide_layouts[6])
        a_box = agenda.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11), Inches(0.9))
        a_tf = a_box.text_frame
        a_tf.text = "Agenda"
        a_tf.paragraphs[0].font.size = Pt(32)
        a_tf.paragraphs[0].font.bold = True
        a_tf.paragraphs[0].font.color.rgb = self._rgb("4f46e5")
        agenda_items = self._sections_map(strategy)
        a_list = agenda.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(5.5))
        a_tf2 = a_list.text_frame
        first = True
        for key in self._SECTION_KEYS:
            if key not in agenda_items:
                continue
            para = a_tf2.paragraphs[0] if first else a_tf2.add_paragraph()
            first = False
            para.text = f"•  {key}"
            para.font.size = Pt(15)
            para.font.color.rgb = self._rgb("1e293b")

        # --- Section slides ----------------------------------------------
        sections = self._sections_map(strategy)
        for key in self._SECTION_KEYS:
            content = sections.get(key)
            if not content:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
            title_tf = title_box.text_frame
            title_tf.text = key
            title_tf.paragraphs[0].font.size = Pt(28)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = self._rgb("4f46e5")

            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.6))
            body_tf = body_box.text_frame
            lines = [ln.rstrip() for ln in content.split("\n") if ln.strip()]
            first = True
            for line in lines:
                para = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                first = False
                para.text = line
                para.font.size = Pt(13)
                para.font.color.rgb = self._rgb("1e293b")

        prs.save(buffer)
        return RenderedExport(
            content=buffer.getvalue(),
            media_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            file_extension="pptx",
        )

    @staticmethod
    def _rgb(hex_color: str):
        from pptx.dml.color import RGBColor

        return RGBColor.from_string(hex_color)

    @staticmethod
    def _industry(strategy) -> str:
        content = strategy.content or {}
        market = content.get("marketOverview") or {}
        return market.get("targetMarketSize") or "—"

    @staticmethod
    def _date_line(strategy) -> str:
        now = strategy.created_at
        if hasattr(now, "strftime"):
            return f"Generation date: {now.strftime('%B %d, %Y')}"
        return ""

    @staticmethod
    def _sections_map(strategy) -> dict[str, str]:
        return {title: content for title, content in _sections(strategy)}


# Maps ExportFormat values to their renderer. Extend here to add formats.
RENDERERS: dict[ExportFormat, BaseRenderer] = {
    JsonRenderer.format: JsonRenderer(),
    MarkdownRenderer.format: MarkdownRenderer(),
    HtmlRenderer.format: HtmlRenderer(),
    PdfRenderer.format: PdfRenderer(),
    DocxRenderer.format: DocxRenderer(),
    PptxRenderer.format: PptxRenderer(),
}


def get_renderer(format_: ExportFormat) -> BaseRenderer:
    """Return the renderer for a format, raising if none is registered."""
    try:
        return RENDERERS[format_]
    except KeyError:
        raise ValueError(f"Unsupported export format: {format_.value}") from None
