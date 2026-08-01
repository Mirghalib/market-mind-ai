"""Professional PowerPoint deck builder.

Builds a consulting-grade, 15+ slide marketing strategy deck from a
MarketingStrategy using python-pptx:

  - Cover slide with brand gradient + meta
  - Agenda
  - Executive summary + marketing score (score bar)
  - Section slides (bulleted, structured)
  - Timeline (roadmap phases)
  - Charts (ROI projection bar chart)
  - Scorecard / KPI tables
  - Final recommendations + closing

All drawing helpers fail soft: missing content simply produces fewer
slides rather than an HTTP 500.
"""
from __future__ import annotations

import logging
from typing import Any

from app.models.marketing_strategy import MarketingStrategy

logger = logging.getLogger("market_mind_ai.pptx")

BRAND = {
    "primary": "4F46E5",   # indigo-600
    "secondary": "7C3AED",  # violet-600
    "accent": "06B6D4",     # cyan-500
    "dark": "0F172A",       # slate-900
    "muted": "64748B",      # slate-500
    "light": "F8FAFC",      # slate-50
    "white": "FFFFFF",
    "emerald": "10B981",
    "amber": "F59E0B",
    "rose": "F43F5E",
}


def _rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color)


def _sections_map(strategy: MarketingStrategy) -> dict[str, str]:
    """Return {title: content} from the stored payload in all shapes.

    Prefers the flattened renderers._sections() titles (human-readable)
    but falls back to compact keys when the payload uses the raw LLM
    document shape.
    """
    from app.services.export.renderers import _sections

    mapping: dict[str, str] = {}
    for title, content in _sections(strategy):
        mapping[title] = content
        # Also index by a normalized compact key (no spaces, lowercase)
        # so structured lookups can find it either way.
        compact = title.replace(" ", "").replace("&", "").replace("(", "").replace(")", "").lower()
        mapping.setdefault(compact, content)
    return mapping


def _find_section(sections: dict[str, str], *names: str) -> str | None:
    """Look up a section by any of the given names (pretty or compact)."""
    for name in names:
        if name in sections:
            return sections[name]
        compact = name.replace(" ", "").replace("&", "").replace("(", "").replace(")", "").lower()
        if compact in sections:
            return sections[compact]
    return None


def _content(strategy: MarketingStrategy) -> dict:
    return strategy.content or {}


def _industry(strategy: MarketingStrategy) -> str:
    market = _content(strategy).get("marketOverview") or {}
    if isinstance(market, dict):
        return market.get("targetMarketSize") or "—"
    return "—"


def _date_line(strategy: MarketingStrategy) -> str:
    now = strategy.created_at
    if hasattr(now, "strftime"):
        return f"Generated {now.strftime('%B %d, %Y')}"
    return ""


class DeckBuilder:
    """Builds the full deck with chart/timeline/table helpers."""

    def __init__(self, strategy: MarketingStrategy) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        self.strategy = strategy
        self.sections = _sections_map(strategy)
        self.content = _content(strategy)
        self.Inches = Inches
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._blank = self.prs.slide_layouts[6]

    # --- shape helpers ------------------------------------------------------

    def add_slide(self):
        return self.prs.slides.add_slide(self._blank)

    def add_rect(self, slide, x, y, w, h, fill, line=None):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, self.Inches(x), self.Inches(y), self.Inches(w), self.Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = self.Inches(0.01)
        shape.shadow.inherit = False
        return shape

    def add_text(self, slide, x, y, w, h, lines, *, align=None, wrap=True):
        """Add a text box. ``lines`` is a list of (text, size, bold, color)."""
        from pptx.enum.text import PP_ALIGN

        box = slide.shapes.add_textbox(
            self.Inches(x), self.Inches(y), self.Inches(w), self.Inches(h)
        )
        tf = box.text_frame
        tf.word_wrap = wrap
        for i, (text, size, bold, color) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = self.Inches(size / 72) if False else _pt(size)
            p.font.bold = bold
            p.font.color.rgb = _rgb(color)
            p.font.name = "Arial"
            if align:
                p.alignment = PP_ALIGN[align]
        return box

    def add_bullets(self, slide, x, y, w, h, items, *, size=14, color="334155"):
        box = slide.shapes.add_textbox(
            self.Inches(x), self.Inches(y), self.Inches(w), self.Inches(h)
        )
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = _pt(size)
            p.font.color.rgb = _rgb(color)
            p.font.name = "Arial"
            p.space_after = _pt(4)
        return box

    def add_kicker(self, slide, x, y, text, color="4F46E5"):
        return self.add_text(
            slide, x, y, 8, 0.4,
            [(text.upper(), 12, True, color)],
        )

    def slide_header(self, slide, title: str, subtitle: str = ""):
        self.add_rect(slide, 0, 0, 13.333, 0.12, BRAND["primary"])
        self.add_text(
            slide, 0.7, 0.35, 11, 0.8,
            [(title, 28, True, BRAND["dark"])],
        )
        if subtitle:
            self.add_text(
                slide, 0.7, 1.05, 11, 0.4,
                [(subtitle, 13, False, BRAND["muted"])],
            )

    def footer(self, slide, page_num: int, total: int):
        self.add_text(
            slide, 0.7, 7.05, 8, 0.3,
            [(f"Market Mind AI  ·  {self.strategy.name}", 9, False, "94A3B8")],
        )
        self.add_text(
            slide, 11.9, 7.05, 0.9, 0.3,
            [(f"{page_num} / {total}", 9, False, "94A3B8")],
            align="RIGHT",
        )

    # --- data helpers -------------------------------------------------------

    def _flat_lines(self, content: str, max_lines: int = 14) -> list[str]:
        """Flatten a section into a list of readable lines."""
        lines = [ln.strip() for ln in (content or "").split("\n") if ln.strip()]
        if not lines:
            return ["No content provided for this section."]
        out: list[str] = []
        for ln in lines:
            out.append(ln)
            if len(out) >= max_lines:
                break
        return out

    def _nested_blocks(self, content: str, max_blocks: int = 8) -> list[tuple[str, list[str]]]:
        """Split a section into (heading, bullet list) blocks.

        Detects "Label:" prefixes; plain lines become one block.
        """
        lines = self._flat_lines(content, max_lines=200)
        blocks: list[tuple[str, list[str]]] = []
        current: list[str] = []
        current_label = ""

        def flush() -> None:
            nonlocal current
            if current:
                blocks.append((current_label, current))
                current = []

        for ln in lines:
            if ":" in ln and len(ln.split(":", 1)[0]) < 28:
                flush()
                head, _, rest = ln.partition(":")
                current_label = head.strip().strip("•").strip()
                if rest.strip():
                    current = [rest.strip()]
                else:
                    current = []
            else:
                current.append(ln.lstrip("• ").strip())
        flush()
        if not blocks:
            blocks = [("", lines)]
        return blocks[:max_blocks]

    # --- slides -------------------------------------------------------------

    def cover(self) -> int:
        from app.services.export.report_data import ReportData

        data = ReportData(self.strategy)
        slide = self.add_slide()
        self.add_rect(slide, 0, 0, 13.333, 7.5, "111827")
        self.add_rect(slide, 0, 0, 13.333, 0.35, BRAND["primary"])
        self.add_text(
            slide, 0.9, 2.0, 11.5, 1.6,
            [(self.strategy.name, 44, True, BRAND["white"])],
        )
        self.add_text(
            slide, 0.9, 3.4, 11.5, 0.6,
            [("Marketing Strategy Report", 24, False, "C7D2FE")],
        )
        budget_line = (
            f"Budget: {data.budget_label}"
            if data.budget_label and data.budget_label != "Not specified"
            else None
        )
        cover_meta = [
            f"Industry: {_industry(self.strategy)}",
            f"Country: {data.country}",
            f"Target audience: {self.strategy.target_audience or '—'}",
        ]
        if budget_line:
            cover_meta.append(budget_line)
        cover_meta.append(_date_line(self.strategy))
        self.add_text(
            slide, 0.9, 4.6, 11.5, 1.6,
            [(line, 14, False, "94A3B8") for line in cover_meta],
        )
        return 1

    def agenda(self) -> int:
        slide = self.add_slide()
        self.slide_header(slide, "Agenda")
        keys = [
            ("Executive Summary", "Executivesummary"),
            ("Marketing Score", "Marketingscore"),
            ("Market Overview", "Marketoverview"),
            ("Customer Persona", "Customerpersona"),
            ("SWOT Analysis", "Swotanalysis"),
            ("Competitor Analysis", "Competitoranalysis"),
            ("Marketing Channels", "Marketingchannels"),
            ("SEO Strategy", "Seostrategy"),
            ("Email Marketing Strategy", "Emailmarketing"),
            ("Social Media Strategy", "Socialmediastrategy"),
            ("Google & Meta Ads Plan", "Advertisementideas"),
            ("Content Calendar", "Contentcalendar"),
            ("Implementation Roadmap (90 Days)", "Implementationroadmap"),
            ("Estimated ROI", "Estimatedroi"),
            ("Risks and Mitigation", "Riskmitigation"),
            ("Final Recommendations", "Finalrecommendations"),
        ]
        items = [pretty for pretty, compact in keys if _find_section(self.sections, pretty, compact)]
        half = (len(items) + 1) // 2
        left = items[:half]
        right = items[half:]
        self.add_bullets(slide, 1.0, 1.8, 5.5, 5.0, left, size=16)
        self.add_bullets(slide, 6.8, 1.8, 5.5, 5.0, right, size=16)
        return 1

    def exec_summary(self) -> int:
        content = _find_section(self.sections, "Executive Summary", "Executivesummary")
        if not content:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Executive Summary")
        blocks = self._nested_blocks(content, max_blocks=5)
        y = 1.7
        for label, lines in blocks[:5]:
            if label:
                self.add_text(slide, 0.9, y, 11.5, 0.4, [(label, 15, True, BRAND["primary"])])
                y += 0.5
            for line in lines[:6]:
                self.add_bullets(slide, 0.9, y, 11.5, 0.4, [line], size=13)
                y += 0.35
            y += 0.25
        return 1

    def business_overview(self) -> int:
        """Business overview slide: goals, positioning, key messages."""
        strategy = self.content.get("marketingStrategy") or {}
        goals = strategy.get("objectives") or []
        positioning = strategy.get("positioning")
        messages = strategy.get("keyMessages") or []
        if not goals and not positioning and not messages:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Business Overview")
        y = 1.8
        if positioning:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Positioning", 15, True, BRAND["primary"])])
            y += 0.5
            self.add_text(slide, 0.9, y, 11.5, 0.8, [(str(positioning), 13, False, BRAND["muted"])])
            y += 1.2
        if goals:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Objectives", 15, True, BRAND["primary"])])
            y += 0.5
            self.add_bullets(slide, 0.9, y, 11.5, 3.0, [str(g) for g in goals[:6]], size=13)
            y += min(6.5 - y, 0.4 * len(goals[:6]))
        if messages and y < 6.2:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Key Messages", 15, True, BRAND["secondary"])])
            y += 0.5
            self.add_bullets(slide, 0.9, y, 11.5, 1.5, [str(m) for m in messages[:4]], size=13)
        return 1

    def target_audience(self) -> int:
        """Target audience slide from the persona."""
        persona = self.content.get("customerPersona") or {}
        if not persona:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Target Audience")
        y = 1.8
        profile = [
            ("Name", persona.get("name")),
            ("Age range", persona.get("ageRange")),
            ("Location", persona.get("location")),
            ("Occupation", persona.get("occupation")),
            ("Income level", persona.get("incomeLevel")),
        ]
        # Left column: profile attributes.
        for label, value in profile:
            if value:
                self.add_text(slide, 0.9, y, 2.4, 0.35, [(label.upper(), 10, True, BRAND["muted"])])
                self.add_text(slide, 0.9, y + 0.32, 5.0, 0.35, [(str(value), 14, True, BRAND["dark"])])
                y += 0.9
        # Right column: interests / goals / pain points.
        x = 6.8
        self.add_text(slide, x, 1.8, 5.5, 0.4, [("Interests", 15, True, BRAND["primary"])])
        interests = persona.get("interests") or []
        self.add_bullets(slide, x, 2.3, 5.5, 2.0, [str(i) for i in interests[:5]], size=12)
        y2 = 4.4
        goals = persona.get("goals") or []
        if goals:
            self.add_text(slide, x, y2, 5.5, 0.4, [("Goals", 15, True, BRAND["secondary"])])
            self.add_bullets(slide, x, y2 + 0.5, 5.5, 1.6, [str(g) for g in goals[:4]], size=12)
        summary = persona.get("summary")
        if summary:
            self.add_text(slide, 0.9, 6.2, 11.5, 0.7, [(str(summary), 12, False, BRAND["muted"])])
        return 1

    def marketing_objectives(self) -> int:
        """Marketing objectives slide with numbered goals."""
        strategy = self.content.get("marketingStrategy") or {}
        goals = strategy.get("objectives") or []
        if not goals:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Marketing Objectives")
        y = 1.9
        for i, goal in enumerate(goals[:8], start=1):
            self.add_rect(slide, 0.9, y, 0.5, 0.5, BRAND["primary"])
            self.add_text(slide, 1.0, y + 0.05, 0.4, 0.4, [(str(i), 16, True, BRAND["white"])], align="CENTER")
            self.add_text(slide, 1.7, y + 0.03, 10.8, 0.5, [(str(goal), 15, False, BRAND["dark"])])
            y += 0.68
        return 1

    def content_strategy(self) -> int:
        """Content strategy slide from the content calendar + SEO topics."""
        calendar = self.content.get("contentCalendar") or {}
        seo = self.content.get("seoKeywords") or {}
        topics = seo.get("contentTopics") or []
        schedule = calendar.get("schedule") or []
        if not topics and not schedule and not calendar.get("timeframe"):
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Content Strategy")
        y = 1.8
        tf = calendar.get("timeframe") or calendar.get("cadence")
        if tf:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Plan", 14, True, BRAND["primary"])])
            y += 0.5
            self.add_text(slide, 0.9, y, 11.5, 0.4, [(str(tf), 13, False, BRAND["dark"])])
            y += 0.8
        if topics:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Content Topics", 15, True, BRAND["secondary"])])
            y += 0.5
            for t in topics[:5]:
                title = t.get("title", "")
                kw = t.get("targetKeyword", "")
                stage = t.get("funnelStage", "")
                line = title + (f" — {kw}" if kw else "")
                self.add_bullets(slide, 0.9, y, 11.5, 0.4, [line], size=13)
                y += 0.4
                if y > 6.2:
                    break
        return 1

    def marketing_score(self) -> int:
        score = self.content.get("marketingScore") or {}
        if not score:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Marketing Score")
        overall = score.get("overall")
        self.add_rect(slide, 0.9, 1.6, 4.2, 2.6, "EEF2FF")
        self.add_text(
            slide, 0.9, 1.9, 4.2, 0.5,
            [(f"{overall}", 60, True, BRAND["primary"])],
            align="CENTER",
        )
        self.add_text(
            slide, 0.9, 3.0, 4.2, 0.4,
            [("OVERALL SCORE / 100", 12, True, BRAND["muted"])],
            align="CENTER",
        )
        # Score bar.
        pct = min(max(int(overall or 0), 0), 100) / 100
        bar_y = 4.6
        self.add_rect(slide, 0.9, bar_y, 11.5, 0.3, "E2E8F0")
        self.add_rect(slide, 0.9, bar_y, 11.5 * pct, 0.3, BRAND["primary"])

        breakdown = score.get("breakdown") or []
        y = 1.6
        for b in breakdown[:8]:
            area = b.get("area", "")
            val = b.get("score")
            if area and val is not None:
                self.add_text(slide, 6.0, y, 3.5, 0.35, [(str(area), 13, True, BRAND["dark"])])
                self.add_rect(slide, 9.6, y + 0.04, 3.2, 0.28, "E2E8F0")
                try:
                    width = 3.2 * min(max(int(val), 0), 100) / 100
                except (TypeError, ValueError):
                    width = 0
                if width:
                    self.add_rect(slide, 9.6, y + 0.04, width, 0.28, BRAND["accent"])
                self.add_text(slide, 12.9, y, 0.6, 0.35, [(str(val), 12, True, BRAND["dark"])])
                y += 0.55
        summary = score.get("summary")
        if summary:
            self.add_text(slide, 0.9, 5.3, 11.5, 1.2, [(str(summary), 13, False, BRAND["muted"])])
        return 1

    def section_slide(self, *names: str) -> int:
        title = names[0]
        content = _find_section(self.sections, *names)
        if not content:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, title)
        blocks = self._nested_blocks(content, max_blocks=6)
        y = 1.7
        for label, lines in blocks:
            if y > 6.4:
                break
            if label:
                self.add_text(slide, 0.9, y, 11.5, 0.4, [(label, 15, True, BRAND["primary"])])
                y += 0.5
            for line in lines[:5]:
                self.add_bullets(slide, 0.9, y, 11.5, 0.4, [line], size=13)
                y += 0.34
            y += 0.2
        return 1

    def roadmap(self) -> int:
        roadmap = self.content.get("implementationRoadmap") or {}
        phases = roadmap.get("phases") or []
        if not phases:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Implementation Roadmap (90 Days)")
        # Timeline bar.
        self.add_rect(slide, 0.9, 1.7, 11.5, 0.06, "CBD5E1")
        n = len(phases[:4])
        step = 11.5 / max(n, 1)
        colors = [BRAND["primary"], BRAND["secondary"], BRAND["accent"], BRAND["emerald"]]
        for i, phase in enumerate(phases[:4]):
            x = 0.9 + step * i
            self.add_rect(slide, x, 1.55, 0.35, 0.35, colors[i % len(colors)])
            self.add_text(
                slide, x - 0.2, 2.1, step + 0.4, 0.4,
                [(phase.get("duration", ""), 12, True, BRAND["primary"])],
            )
            y = 2.7
            self.add_text(slide, x - 0.2, y, step + 0.4, 0.4, [(phase.get("name", ""), 15, True, BRAND["dark"])])
            y += 0.5
            for obj in (phase.get("objectives") or [])[:3]:
                self.add_bullets(slide, x - 0.2, y, step + 0.4, 0.4, [str(obj)], size=11)
                y += 0.3
        return 1

    def roi(self) -> int:
        roi = self.content.get("estimatedROI") or {}
        projections = roi.get("projections") or []
        if not projections:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Estimated ROI")
        # Bar chart of ROI % by period.
        chart_data = []
        for p in projections[:6]:
            try:
                text = str(p.get("roiPercent", "0")).replace("%", "").strip()
                val = float(text)
            except (TypeError, ValueError):
                val = 0.0
            chart_data.append((p.get("period", ""), val))
        max_val = max([abs(v) for _, v in chart_data] + [1])
        chart_x = 0.9
        chart_w = 7.0
        chart_h = 4.2
        base_y = 5.3
        bar_w = min(0.9, chart_w / max(len(chart_data), 1) - 0.3)
        # Axes.
        self.add_rect(slide, chart_x, base_y - chart_h, 0.03, chart_h, "CBD5E1")
        self.add_rect(slide, chart_x, base_y, chart_w, 0.03, "CBD5E1")
        for i, (period, val) in enumerate(chart_data):
            x = chart_x + 0.6 + i * ((chart_w - 0.9) / max(len(chart_data), 1))
            height = max((abs(val) / max_val) * chart_h * 0.85, 0.05)
            color = BRAND["emerald"] if val >= 0 else BRAND["rose"]
            bar_y = base_y - height
            self.add_rect(slide, x, bar_y, bar_w, height, color)
            self.add_text(slide, x - 0.2, bar_y - 0.4, bar_w + 0.4, 0.35, [(f"{val:+.0f}%", 11, True, BRAND["dark"])], align="CENTER")
            self.add_text(slide, x - 0.2, base_y + 0.12, bar_w + 0.4, 0.35, [(str(period), 11, False, BRAND["muted"])], align="CENTER")
        # Summary panel.
        self.add_rect(slide, 8.4, 1.8, 4.0, 4.6, "F8FAFC")
        self.add_text(slide, 8.7, 2.0, 3.4, 0.4, [("ROI Summary", 15, True, BRAND["primary"])])
        y = 2.6
        summary = roi.get("summary")
        if summary:
            self.add_text(slide, 8.7, y, 3.4, 1.4, [(str(summary), 12, False, BRAND["muted"])], wrap=True)
            y += 1.6
        payback = roi.get("paybackPeriod")
        if payback:
            self.add_text(slide, 8.7, y, 3.4, 0.4, [(f"Payback period: {payback}", 13, True, BRAND["dark"])])
            y += 0.5
        for p in projections[:3]:
            self.add_bullets(
                slide, 8.7, y, 3.4, 0.4,
                [f"{p.get('period')}: {p.get('projectedReturn')}"],
                size=12,
            )
            y += 0.35
        return 1

    def recommendations(self) -> int:
        rec = self.content.get("finalRecommendations") or {}
        if not rec:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Final Recommendations")
        summary = rec.get("summary")
        if summary:
            self.add_text(slide, 0.9, 1.7, 11.5, 1.0, [(str(summary), 14, False, BRAND["muted"])])
        y = 2.9
        for label, key, color in [
            ("Top priorities", "priorities", BRAND["primary"]),
            ("Quick wins", "quickWins", BRAND["emerald"]),
            ("Long-term investments", "longTermInvestments", BRAND["secondary"]),
        ]:
            items = rec.get(key) or []
            if not items:
                continue
            self.add_text(slide, 0.9, y, 11.5, 0.4, [(label, 15, True, color)])
            y += 0.5
            for item in items[:4]:
                self.add_bullets(slide, 0.9, y, 11.5, 0.4, [str(item)], size=13)
                y += 0.34
            y += 0.2
            if y > 6.4:
                break
        return 1

    def kpi_slide(self) -> int:
        """KPIs / scorecard table from the strategy content."""
        strategy = self.content.get("marketingStrategy") or {}
        kpis = strategy.get("kpis") or []
        if not kpis:
            # Fall back to a KPI summary from the flattened section.
            content = _find_section(self.sections, "KPIs", "Kpis", "Kpisummary")
            if not content:
                return 0
            slide = self.add_slide()
            self.slide_header(slide, "Key Performance Indicators")
            blocks = self._nested_blocks(content, max_blocks=6)
            y = 1.8
            for label, lines in blocks:
                if y > 6.4:
                    break
                if label:
                    self.add_text(slide, 0.9, y, 11.5, 0.4, [(label, 15, True, BRAND["primary"])])
                    y += 0.5
                for line in lines[:5]:
                    self.add_bullets(slide, 0.9, y, 11.5, 0.4, [line], size=13)
                    y += 0.34
                y += 0.2
            return 1

        slide = self.add_slide()
        self.slide_header(slide, "Key Performance Indicators")
        # Table header.
        cols = [(0.9, 4.5, "Metric"), (5.6, 2.5, "Target"), (8.2, 2.0, "Timeframe"), (10.3, 2.0, "Owner")]
        for x, w, label in cols:
            self.add_rect(slide, x, 1.7, w, 0.5, BRAND["primary"])
            self.add_text(slide, x + 0.15, 1.78, w - 0.3, 0.35, [(label, 12, True, BRAND["white"])])
        y = 2.25
        for i, kpi in enumerate(kpis[:8]):
            bg = "F8FAFC" if i % 2 == 0 else "FFFFFF"
            for x, w, _ in cols:
                self.add_rect(slide, x, y, w, 0.55, bg)
            self.add_text(slide, 0.9 + 0.15, y + 0.1, 4.2, 0.35, [(str(kpi.get("metric", "")), 12, True, BRAND["dark"])])
            self.add_text(slide, 5.6 + 0.15, y + 0.1, 2.2, 0.35, [(str(kpi.get("target", "")), 12, False, BRAND["emerald"])])
            self.add_text(slide, 8.2 + 0.15, y + 0.1, 1.7, 0.35, [(str(kpi.get("timeframe", "")), 11, False, BRAND["muted"])])
            y += 0.55
        return 1

    def channels_slide(self) -> int:
        """Marketing channels + budget allocation slide."""
        strategy = self.content.get("marketingStrategy") or {}
        channels = strategy.get("channels") or []
        allocation = strategy.get("budgetAllocation") or []
        if not channels and not allocation:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Marketing Channels & Budget")
        y = 1.8
        if channels:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Channels", 15, True, BRAND["primary"])])
            y += 0.5
            for ch in channels[:6]:
                self.add_bullets(
                    slide, 0.9, y, 11.5, 0.4,
                    [f"{ch.get('name', '')} ({ch.get('priority', '')}): {ch.get('description', '')}"],
                    size=13,
                )
                y += 0.38
            y += 0.3
        if allocation:
            self.add_text(slide, 0.9, y, 11.5, 0.4, [("Budget allocation", 15, True, BRAND["secondary"])])
            y += 0.5
            for a in allocation[:6]:
                try:
                    pct = float(a.get("percentage", 0))
                except (TypeError, ValueError):
                    pct = 0
                self.add_text(slide, 0.9, y, 4.5, 0.35, [(str(a.get("channel", "")), 13, True, BRAND["dark"])])
                self.add_rect(slide, 5.6, y + 0.04, 6.0, 0.28, "E2E8F0")
                if pct:
                    self.add_rect(slide, 5.6, y + 0.04, min(6.0 * pct / 100, 6.0), 0.28, BRAND["accent"])
                self.add_text(slide, 11.8, y, 0.9, 0.35, [(f"{pct:.0f}%", 12, True, BRAND["dark"])])
                y += 0.5
        return 1

    def next_steps(self) -> int:
        """Next steps slide from the roadmap + recommendations."""
        roadmap = self.content.get("implementationRoadmap") or {}
        rec = self.content.get("finalRecommendations") or {}
        phases = roadmap.get("phases") or []
        steps = [f"{p.get('name', '')} — {p.get('duration', '')}" for p in phases[:4]]
        priorities = rec.get("priorities") or []
        items = [str(p) for p in priorities[:4]] if priorities else steps
        if not items:
            return 0
        slide = self.add_slide()
        self.slide_header(slide, "Next Steps")
        y = 1.8
        for i, item in enumerate(items[:6], start=1):
            self.add_rect(slide, 0.9, y, 0.5, 0.5, BRAND["emerald"])
            self.add_text(slide, 1.0, y + 0.05, 0.4, 0.4, [(str(i), 15, True, BRAND["white"])], align="CENTER")
            self.add_text(slide, 1.7, y + 0.04, 10.8, 0.5, [(str(item), 15, False, BRAND["dark"])])
            y += 0.68
        quick = rec.get("quickWins") or []
        if quick and y < 5.6:
            self.add_text(slide, 0.9, y + 0.2, 11.5, 0.4, [("Quick wins", 14, True, BRAND["secondary"])])
            y += 0.7
            self.add_bullets(slide, 0.9, y, 11.5, 1.2, [str(q) for q in quick[:4]], size=12)
        return 1

    def closing(self) -> int:
        slide = self.add_slide()
        self.add_rect(slide, 0, 0, 13.333, 7.5, "111827")
        self.add_rect(slide, 0, 0, 13.333, 0.35, BRAND["primary"])
        self.add_text(
            slide, 0.9, 2.6, 11.5, 1.0,
            [("Thank you", 44, True, BRAND["white"])],
        )
        self.add_text(
            slide, 0.9, 3.7, 11.5, 1.2,
            [
                (
                    "This strategy was generated by Market Mind AI — AI-powered "
                    "marketing intelligence.",
                    16, False, "94A3B8",
                ),
                (f"{self.strategy.name}  ·  {_date_line(self.strategy)}", 13, False, "64748B"),
            ],
        )
        return 1

    # --- build --------------------------------------------------------------

    def build(self) -> bytes:
        from io import BytesIO

        order = [
            self.cover,
            self.agenda,
            self.exec_summary,
            self.marketing_score,
            self.business_overview,
            self.target_audience,
            self.marketing_objectives,
            self.kpi_slide,
            lambda: self.section_slide("Market Overview"),
            lambda: self.section_slide("Customer Persona"),
            lambda: self.section_slide("SWOT Analysis"),
            lambda: self.section_slide("Competitor Analysis"),
            self.channels_slide,
            self.content_strategy,
            lambda: self.section_slide("SEO Strategy", "Seokeywords"),
            lambda: self.section_slide("Email Marketing Strategy", "Emailcampaign"),
            lambda: self.section_slide("Social Media Strategy", "Socialmediastrategy"),
            lambda: self.section_slide("Google & Meta Ads Plan", "Advertisementideas"),
            lambda: self.section_slide("Content Calendar", "Contentcalendar"),
            self.roadmap,
            self.roi,
            lambda: self.section_slide("Risks and Mitigation", "Riskmitigation"),
            self.recommendations,
            self.next_steps,
            self.closing,
        ]

        for build_fn in order:
            try:
                build_fn()
            except Exception:  # noqa: BLE001
                logger.exception("PPTX slide builder failed for %s", getattr(build_fn, "__name__", "slide"))
                continue

        # Number slides.
        total = len(self.prs.slides._sldIdLst)
        for idx, slide in enumerate(self.prs.slides, start=1):
            try:
                if idx > 1 and idx < total:
                    self.footer(slide, idx, total)
            except Exception:  # noqa: BLE001
                continue

        buffer = BytesIO()
        self.prs.save(buffer)
        return buffer.getvalue()


def build_deck(strategy: MarketingStrategy) -> bytes:
    """Build the professional deck and return its bytes."""
    return DeckBuilder(strategy).build()


def _pt(size: int | float):
    from pptx.util import Pt

    return Pt(size)
